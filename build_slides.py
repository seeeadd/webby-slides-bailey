"""
Bailey Vann - The 2026 Etsy Reset
Premium Slide Deck Builder
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.enum.dml import MSO_THEME_COLOR
import os

# =============================================================================
# DESIGN SYSTEM
# =============================================================================

class Colors:
    """Bailey Vann Brand Colors"""
    # Primary
    TEAL_DEEP = RGBColor(0x1B, 0x8A, 0x8A)      # #1B8A8A - authority, headers
    TEAL_LIGHT = RGBColor(0x2B, 0xA5, 0xA3)     # #2BA5A3 - lighter teal
    CORAL = RGBColor(0xE0, 0x7B, 0x6C)          # #E07B6C - emphasis, warmth
    CORAL_SOFT = RGBColor(0xF4, 0xA8, 0x9A)     # #F4A89A - gentler highlights

    # Backgrounds
    CREAM = RGBColor(0xFD, 0xF8, 0xF3)          # #FDF8F3 - primary bg
    BLUSH = RGBColor(0xFE, 0xE5, 0xE0)          # #FEE5E0 - section variation
    MINT_SOFT = RGBColor(0xE8, 0xF5, 0xF3)      # #E8F5F3 - subtle contrast
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # #FFFFFF

    # Accents
    GOLD = RGBColor(0xD4, 0xAF, 0x37)           # #D4AF37 - premium accent
    YELLOW_WARM = RGBColor(0xF7, 0xD7, 0x6C)    # #F7D76C - sparkle

    # Text
    TEXT_DARK = RGBColor(0x2D, 0x34, 0x36)      # #2D3436 - primary text
    TEXT_MUTED = RGBColor(0x63, 0x6E, 0x72)     # #636E72 - secondary text
    TEXT_LIGHT = RGBColor(0x9C, 0xA3, 0xA8)     # #9CA3A8 - subtle text


class Fonts:
    """Typography System"""
    # Headers - OGG (elegant serif) - actual font family name from files
    HEADER = "Ogg TRIAL"
    HEADER_MEDIUM = "Ogg TRIAL Medium"
    HEADER_LIGHT = "Ogg TRIAL Light"

    # Ogg Text variants (for smaller text)
    TEXT = "Ogg Text TRIAL"
    TEXT_BOLD = "Ogg Text TRIAL"
    TEXT_MEDIUM = "Ogg Text TRIAL Medium"

    # Body - Arial as reliable fallback (Satoshi needs to be installed)
    BODY = "Arial"

    # Mono - for data/numbers
    MONO = "Consolas"

    # Script - for annotations
    SCRIPT = "Brush Script MT"


class Sizing:
    """Consistent sizing"""
    # Slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Margins
    MARGIN = Inches(0.75)
    MARGIN_WIDE = Inches(1.0)

    # Font sizes
    TITLE_HUGE = Pt(72)
    TITLE_LARGE = Pt(54)
    TITLE_MEDIUM = Pt(42)
    HEADING = Pt(32)
    SUBHEADING = Pt(24)
    BODY_LARGE = Pt(20)
    BODY = Pt(18)
    BODY_SMALL = Pt(14)
    CAPTION = Pt(12)

    # Corner radius for shapes (in EMUs)
    CORNER_RADIUS = Emu(150000)  # ~12px


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def create_presentation():
    """Create a new presentation with 16:9 aspect ratio"""
    prs = Presentation()
    prs.slide_width = Sizing.SLIDE_WIDTH
    prs.slide_height = Sizing.SLIDE_HEIGHT
    return prs


def add_blank_slide(prs):
    """Add a blank slide"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)


def set_background_color(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_gradient_background(slide, color1, color2, angle=90):
    """Set gradient background (top to bottom by default)"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2


def add_text_box(slide, left, top, width, height, text,
                 font_name=Fonts.BODY, font_size=Sizing.BODY,
                 font_color=Colors.TEXT_DARK, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP):
    """Add a text box with specified styling"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment

    tf.vertical_anchor = vertical

    return txBox


def add_shape(slide, shape_type, left, top, width, height,
              fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a shape with styling"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()

    return shape


def add_rounded_rectangle(slide, left, top, width, height,
                          fill_color=Colors.WHITE, shadow=True):
    """Add a rounded rectangle (card style)"""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      left, top, width, height, fill_color)

    # Adjust corner radius
    shape.adjustments[0] = 0.05  # 5% corner radius

    # Add shadow if requested
    if shadow:
        add_shadow(shape)

    return shape


def add_shadow(shape):
    """Add a soft shadow to a shape"""
    # Access shadow through XML manipulation
    spPr = shape._sp.spPr

    shadow_xml = '''
    <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="76200" dist="38100" dir="5400000" algn="tl" rotWithShape="0">
            <a:srgbClr val="000000">
                <a:alpha val="15000"/>
            </a:srgbClr>
        </a:outerShdw>
    </a:effectLst>
    '''
    effect_lst = parse_xml(shadow_xml)
    spPr.append(effect_lst)


def add_accent_bar(slide, position='top', color=Colors.TEAL_DEEP, height=Inches(0.12)):
    """Add a thin accent bar at top or bottom of slide"""
    if position == 'top':
        top = Inches(0)
    else:
        top = Sizing.SLIDE_HEIGHT - height

    bar = add_shape(slide, MSO_SHAPE.RECTANGLE,
                    Inches(0), top, Sizing.SLIDE_WIDTH, height, color)
    return bar


def add_pill_badge(slide, left, top, text, bg_color=Colors.CORAL, text_color=Colors.WHITE):
    """Add a pill-shaped badge/label"""
    # Calculate width based on text length
    width = Inches(len(text) * 0.12 + 0.4)
    height = Inches(0.35)

    pill = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                     left, top, width, height, bg_color)
    pill.adjustments[0] = 0.5  # Maximum roundness for pill shape

    # Add text
    tf = pill.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = Fonts.BODY
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    return pill


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def build_slide_1_title(prs):
    """
    SLIDE 1: Title Slide - THE 2026 ETSY RESET

    Design: Bold, editorial, magazine-cover feeling
    Layout: Asymmetric with massive "RESET" as hero
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    # Top accent bar - teal
    add_accent_bar(slide, 'top', Colors.TEAL_DEEP, Inches(0.08))

    # Bottom accent bar - coral
    add_accent_bar(slide, 'bottom', Colors.CORAL, Inches(0.08))

    # Day indicator badge
    add_pill_badge(slide, Inches(0.75), Inches(0.5),
                   "DAY 1 OF 3", Colors.TEAL_DEEP, Colors.WHITE)

    # Main title area - left aligned for editorial feel

    # "THE 2026 ETSY" - smaller, elegant
    add_text_box(slide,
                 Inches(0.75), Inches(1.8), Inches(10), Inches(0.6),
                 "THE 2026 ETSY",
                 font_name=Fonts.HEADER, font_size=Pt(28),
                 font_color=Colors.TEAL_DEEP, bold=False,
                 alignment=PP_ALIGN.LEFT)

    # "RESET" - MASSIVE, the hero word
    add_text_box(slide,
                 Inches(0.65), Inches(2.2), Inches(12), Inches(1.8),
                 "RESET",
                 font_name=Fonts.HEADER, font_size=Pt(140),
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # Subtitle
    add_text_box(slide,
                 Inches(0.75), Inches(4.1), Inches(8), Inches(0.8),
                 "Delete the Dead Weight. Find Your Focus.\nBuild a Shop That Actually Works.",
                 font_name=Fonts.BODY, font_size=Pt(22),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.LEFT)

    # Bailey's credentials - bottom card
    cred_card = add_rounded_rectangle(slide,
                                       Inches(0.75), Inches(5.8),
                                       Inches(7.5), Inches(0.9),
                                       Colors.TEAL_DEEP, shadow=False)

    # Credentials text
    cred_tf = cred_card.text_frame
    cred_tf.clear()
    p = cred_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run1 = p.add_run()
    run1.text = "with "
    run1.font.name = Fonts.BODY
    run1.font.size = Pt(14)
    run1.font.color.rgb = Colors.WHITE

    run2 = p.add_run()
    run2.text = "Bailey Vann"
    run2.font.name = Fonts.BODY
    run2.font.size = Pt(14)
    run2.font.color.rgb = Colors.WHITE
    run2.font.bold = True

    run3 = p.add_run()
    run3.text = "  •  Top 0.1% Etsy Seller  •  $1M+ in Digital Product Sales"
    run3.font.name = Fonts.BODY
    run3.font.size = Pt(14)
    run3.font.color.rgb = Colors.WHITE

    cred_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Decorative element - abstract shape on right side
    # Large teal circle (partial, extends off slide)
    circle = add_shape(slide, MSO_SHAPE.OVAL,
                       Inches(9.5), Inches(1.5),
                       Inches(5), Inches(5),
                       Colors.TEAL_LIGHT)
    circle.fill.solid()
    circle.fill.fore_color.rgb = Colors.TEAL_LIGHT
    # Make it semi-transparent

    # Smaller coral accent circle
    circle2 = add_shape(slide, MSO_SHAPE.OVAL,
                        Inches(10.8), Inches(4.5),
                        Inches(2.5), Inches(2.5),
                        Colors.CORAL_SOFT)

    return slide


def build_slide_2_before_we_begin(prs):
    """
    SLIDE 2: Before We Begin - Quiz intro

    Design: Clean, intriguing, sets up engagement
    No emojis - use visual design instead
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    # Accent bars
    add_accent_bar(slide, 'top', Colors.TEAL_DEEP, Inches(0.06))
    add_accent_bar(slide, 'bottom', Colors.CORAL, Inches(0.06))

    # Center card for content
    card = add_rounded_rectangle(slide,
                                  Inches(2.5), Inches(1.5),
                                  Inches(8.333), Inches(4.5),
                                  Colors.WHITE, shadow=True)

    # "Before We Begin..." - main headline
    add_text_box(slide,
                 Inches(2.5), Inches(2.0), Inches(8.333), Inches(1.0),
                 "Before We Begin...",
                 font_name=Fonts.HEADER, font_size=Sizing.TITLE_MEDIUM,
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Visual divider line
    divider = add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(5.5), Inches(3.1), Inches(2.333), Inches(0.03),
                        Colors.TEAL_DEEP)

    # "Quick Pop Quiz" - with visual badge instead of emoji
    quiz_badge = add_rounded_rectangle(slide,
                                        Inches(4.8), Inches(3.5),
                                        Inches(3.7), Inches(0.6),
                                        Colors.CORAL, shadow=False)

    quiz_tf = quiz_badge.text_frame
    quiz_tf.clear()
    p = quiz_tf.paragraphs[0]
    p.text = "QUICK POP QUIZ"
    p.font.name = Fonts.BODY
    p.font.size = Pt(16)
    p.font.color.rgb = Colors.WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    quiz_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Subtext
    add_text_box(slide,
                 Inches(2.5), Inches(4.5), Inches(8.333), Inches(1.0),
                 "I want to show you something that might change\nhow you think about Etsy in 2026...",
                 font_name=Fonts.BODY, font_size=Pt(18),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_3_get_ready_chat(prs):
    """
    SLIDE 3: Get ready to type in chat

    Design: Engaging, interactive feeling
    Custom chat bubble graphic instead of emoji
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.BLUSH)

    # Accent bars
    add_accent_bar(slide, 'top', Colors.CORAL, Inches(0.06))

    # Large chat bubble graphic - custom built
    # Main bubble
    bubble = add_rounded_rectangle(slide,
                                    Inches(3.5), Inches(2.0),
                                    Inches(6.333), Inches(2.5),
                                    Colors.WHITE, shadow=True)

    # Bubble tail (triangle pointing down-left)
    tail = add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE,
                     Inches(4.2), Inches(4.4),
                     Inches(0.6), Inches(0.5),
                     Colors.WHITE)
    tail.rotation = 180

    # Text inside bubble
    add_text_box(slide,
                 Inches(3.8), Inches(2.5), Inches(5.733), Inches(1.5),
                 "Get ready to type\nin the chat!",
                 font_name=Fonts.HEADER, font_size=Pt(36),
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Decorative typing indicator dots
    for i, offset in enumerate([0, 0.4, 0.8]):
        dot = add_shape(slide, MSO_SHAPE.OVAL,
                        Inches(6.0 + offset), Inches(5.5),
                        Inches(0.25), Inches(0.25),
                        Colors.TEAL_DEEP)

    # Instruction text below
    add_text_box(slide,
                 Inches(2), Inches(6.0), Inches(9.333), Inches(0.6),
                 "This is interactive — your answers matter",
                 font_name=Fonts.BODY, font_size=Pt(16),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_4_quiz_ab(prs):
    """
    SLIDE 4: Which design was made by professional artist? A/B

    Design: Clean comparison layout, two distinct options
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.TEAL_DEEP, Inches(0.06))

    # Question header
    add_text_box(slide,
                 Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
                 "Which design was made by a professional artist?",
                 font_name=Fonts.HEADER, font_size=Pt(32),
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Two side-by-side cards for A and B
    # Card A
    card_a = add_rounded_rectangle(slide,
                                    Inches(0.75), Inches(1.6),
                                    Inches(5.5), Inches(4.8),
                                    Colors.WHITE, shadow=True)

    # Image placeholder A
    img_placeholder_a = add_rounded_rectangle(slide,
                                               Inches(1.0), Inches(1.9),
                                               Inches(5.0), Inches(3.5),
                                               Colors.MINT_SOFT, shadow=False)

    add_text_box(slide,
                 Inches(1.0), Inches(3.2), Inches(5.0), Inches(0.8),
                 "[Image A]",
                 font_name=Fonts.BODY, font_size=Pt(16),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER,
                 vertical=MSO_ANCHOR.MIDDLE)

    # "A" label
    label_a = add_shape(slide, MSO_SHAPE.OVAL,
                        Inches(3.0), Inches(5.6),
                        Inches(0.7), Inches(0.7),
                        Colors.TEAL_DEEP)

    a_tf = label_a.text_frame
    a_tf.clear()
    p = a_tf.paragraphs[0]
    p.text = "A"
    p.font.name = Fonts.HEADER
    p.font.size = Pt(24)
    p.font.color.rgb = Colors.WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    a_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Card B
    card_b = add_rounded_rectangle(slide,
                                    Inches(7.083), Inches(1.6),
                                    Inches(5.5), Inches(4.8),
                                    Colors.WHITE, shadow=True)

    # Image placeholder B
    img_placeholder_b = add_rounded_rectangle(slide,
                                               Inches(7.333), Inches(1.9),
                                               Inches(5.0), Inches(3.5),
                                               Colors.MINT_SOFT, shadow=False)

    add_text_box(slide,
                 Inches(7.333), Inches(3.2), Inches(5.0), Inches(0.8),
                 "[Image B]",
                 font_name=Fonts.BODY, font_size=Pt(16),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER,
                 vertical=MSO_ANCHOR.MIDDLE)

    # "B" label
    label_b = add_shape(slide, MSO_SHAPE.OVAL,
                        Inches(9.333), Inches(5.6),
                        Inches(0.7), Inches(0.7),
                        Colors.CORAL)

    b_tf = label_b.text_frame
    b_tf.clear()
    p = b_tf.paragraphs[0]
    p.text = "B"
    p.font.name = Fonts.HEADER
    p.font.size = Pt(24)
    p.font.color.rgb = Colors.WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    b_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # "VS" in the middle
    add_text_box(slide,
                 Inches(6.0), Inches(3.5), Inches(1.333), Inches(0.6),
                 "VS",
                 font_name=Fonts.HEADER, font_size=Pt(20),
                 font_color=Colors.TEXT_MUTED, bold=True,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_5_type_ab(prs):
    """
    SLIDE 5: Type A or B in the chat

    Design: High energy, call to action
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.TEAL_DEEP)

    # Main message - white on teal
    add_text_box(slide,
                 Inches(1), Inches(2.5), Inches(11.333), Inches(1.2),
                 "Type A or B in the chat!",
                 font_name=Fonts.HEADER, font_size=Pt(54),
                 font_color=Colors.WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Two option buttons side by side
    # A button
    btn_a = add_rounded_rectangle(slide,
                                   Inches(4.0), Inches(4.2),
                                   Inches(2.0), Inches(1.5),
                                   Colors.WHITE, shadow=True)

    a_tf = btn_a.text_frame
    a_tf.clear()
    p = a_tf.paragraphs[0]
    p.text = "A"
    p.font.name = Fonts.HEADER
    p.font.size = Pt(48)
    p.font.color.rgb = Colors.TEAL_DEEP
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    a_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # B button
    btn_b = add_rounded_rectangle(slide,
                                   Inches(7.333), Inches(4.2),
                                   Inches(2.0), Inches(1.5),
                                   Colors.CORAL, shadow=True)

    b_tf = btn_b.text_frame
    b_tf.clear()
    p = b_tf.paragraphs[0]
    p.text = "B"
    p.font.name = Fonts.HEADER
    p.font.size = Pt(48)
    p.font.color.rgb = Colors.WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    b_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # "or" between them
    add_text_box(slide,
                 Inches(6.0), Inches(4.5), Inches(1.333), Inches(0.7),
                 "or",
                 font_name=Fonts.BODY, font_size=Pt(24),
                 font_color=Colors.WHITE,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_6_answer_is(prs):
    """
    SLIDE 6: The Answer Is...

    Design: Suspenseful, dramatic pause
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.TEAL_DEEP, Inches(0.06))
    add_accent_bar(slide, 'bottom', Colors.CORAL, Inches(0.06))

    # Dramatic centered text
    add_text_box(slide,
                 Inches(1), Inches(2.8), Inches(11.333), Inches(1.5),
                 "The Answer Is...",
                 font_name=Fonts.HEADER, font_size=Pt(64),
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtle dots indicating suspense
    for i in range(3):
        dot = add_shape(slide, MSO_SHAPE.OVAL,
                        Inches(6.2 + i * 0.35), Inches(4.5),
                        Inches(0.15), Inches(0.15),
                        Colors.TEAL_DEEP)

    return slide


def build_slide_7_both_ai(prs):
    """
    SLIDE 7: Both Were Made by AI

    Design: REVEAL moment - dramatic, surprising
    Full impact typography
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CORAL)

    # Main reveal - large text
    add_text_box(slide,
                 Inches(0.75), Inches(1.8), Inches(11.833), Inches(1.5),
                 "Both Were Made by AI.",
                 font_name=Fonts.HEADER, font_size=Pt(60),
                 font_color=Colors.WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Supporting text in card
    detail_card = add_rounded_rectangle(slide,
                                         Inches(3.5), Inches(4.0),
                                         Inches(6.333), Inches(1.5),
                                         Colors.WHITE, shadow=True)

    add_text_box(slide,
                 Inches(3.5), Inches(4.3), Inches(6.333), Inches(1.0),
                 "In less than 30 seconds.\nFor free.",
                 font_name=Fonts.BODY, font_size=Pt(24),
                 font_color=Colors.TEXT_DARK, bold=False,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_8_ai_nowadays(prs):
    """
    SLIDE 8: Yes, we all use AI nowadays...

    Design: Contemplative, building tension
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.TEAL_DEEP, Inches(0.06))

    # Main content card
    card = add_rounded_rectangle(slide,
                                  Inches(1.5), Inches(1.5),
                                  Inches(10.333), Inches(4.5),
                                  Colors.WHITE, shadow=True)

    # First part - acknowledgment
    add_text_box(slide,
                 Inches(2.0), Inches(2.0), Inches(9.333), Inches(0.8),
                 "Yes, we all use AI nowadays...",
                 font_name=Fonts.HEADER, font_size=Pt(32),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.LEFT)

    # Divider
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(2.0), Inches(3.0), Inches(3), Inches(0.02),
              Colors.TEAL_DEEP)

    # Key message - emphasized
    add_text_box(slide,
                 Inches(2.0), Inches(3.3), Inches(9.333), Inches(2.0),
                 "But in 2026, they've gotten so good that\nhiring someone REAL is starting to become\njust an option...",
                 font_name=Fonts.BODY, font_size=Pt(24),
                 font_color=Colors.TEXT_DARK,
                 alignment=PP_ALIGN.LEFT)

    # Emphasis on "just an option"
    emphasis_box = add_rounded_rectangle(slide,
                                          Inches(2.0), Inches(5.3),
                                          Inches(2.8), Inches(0.5),
                                          Colors.CORAL_SOFT, shadow=False)

    add_text_box(slide,
                 Inches(2.0), Inches(5.35), Inches(2.8), Inches(0.4),
                 "just an option",
                 font_name=Fonts.BODY, font_size=Pt(16),
                 font_color=Colors.CORAL, bold=True,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_9_sink_in(prs):
    """
    SLIDE 9: Let that sink in for a second.

    Design: Minimal, breathing room, moment of reflection
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.MINT_SOFT)

    # Very minimal - just the text centered
    add_text_box(slide,
                 Inches(1), Inches(3.0), Inches(11.333), Inches(1.5),
                 "Let that sink in for a second.",
                 font_name=Fonts.HEADER, font_size=Pt(42),
                 font_color=Colors.TEAL_DEEP, bold=False,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_10_uncomfortable(prs):
    """
    SLIDE 10: Now let me ask you something uncomfortable...

    Design: Tension building, darker tone
    """
    slide = add_blank_slide(prs)
    # Darker background for tension
    set_background_color(slide, Colors.TEXT_DARK)

    # Main text - white for contrast
    add_text_box(slide,
                 Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.0),
                 "Now let me ask you",
                 font_name=Fonts.HEADER, font_size=Pt(36),
                 font_color=Colors.WHITE,
                 alignment=PP_ALIGN.CENTER)

    # "something uncomfortable" - emphasized
    add_text_box(slide,
                 Inches(1.5), Inches(3.5), Inches(10.333), Inches(1.2),
                 "something uncomfortable...",
                 font_name=Fonts.HEADER, font_size=Pt(48),
                 font_color=Colors.CORAL, bold=True,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_11_what_happens(prs):
    """
    SLIDE 11: If anyone can create designs like this in seconds...
    What happens to YOUR Etsy shop in 2026?

    Design: The BIG question - dramatic, impactful
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.CORAL, Inches(0.08))

    # Setup line
    add_text_box(slide,
                 Inches(1), Inches(1.5), Inches(11.333), Inches(1.0),
                 "If anyone can create designs like this in seconds...",
                 font_name=Fonts.BODY, font_size=Pt(24),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER)

    # THE BIG QUESTION - massive
    add_text_box(slide,
                 Inches(0.75), Inches(2.8), Inches(11.833), Inches(2.5),
                 "What happens to YOUR\nEtsy shop in 2026?",
                 font_name=Fonts.HEADER, font_size=Pt(54),
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Visual emphasis - question mark graphic
    qmark = add_text_box(slide,
                          Inches(10.5), Inches(4.8), Inches(2), Inches(2),
                          "?",
                          font_name=Fonts.HEADER, font_size=Pt(120),
                          font_color=Colors.CORAL_SOFT,
                          alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_12_survey_intro(prs):
    """
    SLIDE 12: The numbers I'm about to show you aren't random...

    Design: Setting up credibility for data
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.TEAL_DEEP, Inches(0.06))

    # Main message
    add_text_box(slide,
                 Inches(1.5), Inches(2.0), Inches(10.333), Inches(1.5),
                 "The numbers I'm about to show you\naren't random...",
                 font_name=Fonts.HEADER, font_size=Pt(40),
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Credibility card
    cred_card = add_rounded_rectangle(slide,
                                       Inches(3), Inches(4.2),
                                       Inches(7.333), Inches(1.8),
                                       Colors.TEAL_DEEP, shadow=True)

    add_text_box(slide,
                 Inches(3.2), Inches(4.5), Inches(6.933), Inches(1.3),
                 "These are from YOUR OWN ANSWERS\nafter I asked 160+ of you last week",
                 font_name=Fonts.BODY, font_size=Pt(20),
                 font_color=Colors.WHITE,
                 alignment=PP_ALIGN.CENTER)

    # "160+" highlighted
    badge_160 = add_rounded_rectangle(slide,
                                       Inches(5.8), Inches(6.2),
                                       Inches(1.7), Inches(0.5),
                                       Colors.CORAL, shadow=False)

    b_tf = badge_160.text_frame
    b_tf.clear()
    p = b_tf.paragraphs[0]
    p.text = "160+ sellers"
    p.font.name = Fonts.BODY
    p.font.size = Pt(12)
    p.font.color.rgb = Colors.WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    b_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide


def build_slide_13_stat_75(prs):
    """
    SLIDE 13: 75+ sellers report shops tanked

    Design: Data visualization - human figures showing proportion
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.CORAL, Inches(0.06))

    # BIG NUMBER
    add_text_box(slide,
                 Inches(0.75), Inches(0.8), Inches(4), Inches(1.8),
                 "75+",
                 font_name=Fonts.HEADER, font_size=Pt(96),
                 font_color=Colors.CORAL, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # Context
    add_text_box(slide,
                 Inches(0.75), Inches(2.5), Inches(5), Inches(1.0),
                 "sellers report their shops\ntanked in the last 2 months",
                 font_name=Fonts.BODY, font_size=Pt(22),
                 font_color=Colors.TEXT_DARK,
                 alignment=PP_ALIGN.LEFT)

    # Visual representation - grid of dots showing 75 out of 160
    # Create a visual grid of small circles
    start_x = Inches(6.5)
    start_y = Inches(1.2)
    dot_size = Inches(0.15)
    gap = Inches(0.22)

    # 10x16 grid = 160 dots, first 75 are coral (affected)
    for row in range(10):
        for col in range(16):
            idx = row * 16 + col
            if idx < 160:
                x = start_x + col * gap
                y = start_y + row * gap
                color = Colors.CORAL if idx < 75 else Colors.MINT_SOFT
                add_shape(slide, MSO_SHAPE.OVAL,
                          x, y, dot_size, dot_size, color)

    # Insight text
    insight_card = add_rounded_rectangle(slide,
                                          Inches(0.75), Inches(4.0),
                                          Inches(5.5), Inches(1.5),
                                          Colors.WHITE, shadow=True)

    add_text_box(slide,
                 Inches(1.0), Inches(4.2), Inches(5.0), Inches(1.2),
                 "If there are 4 people on your row in this chat,\n3 of them are feeling this...",
                 font_name=Fonts.BODY, font_size=Pt(16),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.LEFT)

    return slide


def build_slide_14_stat_58(prs):
    """
    SLIDE 14: 58+ sellers said no views anymore

    Design: Data visualization - flatlined graph
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.TEAL_DEEP, Inches(0.06))

    # BIG NUMBER
    add_text_box(slide,
                 Inches(7.5), Inches(0.8), Inches(5), Inches(1.8),
                 "58+",
                 font_name=Fonts.HEADER, font_size=Pt(96),
                 font_color=Colors.TEAL_DEEP, bold=True,
                 alignment=PP_ALIGN.RIGHT)

    # Context
    add_text_box(slide,
                 Inches(7.5), Inches(2.5), Inches(5), Inches(1.0),
                 "sellers said their #1 problem:\n\"No views anymore\"",
                 font_name=Fonts.BODY, font_size=Pt(22),
                 font_color=Colors.TEXT_DARK,
                 alignment=PP_ALIGN.RIGHT)

    # Flatline graph visualization
    # Graph container
    graph_card = add_rounded_rectangle(slide,
                                        Inches(0.75), Inches(1.5),
                                        Inches(6), Inches(3.5),
                                        Colors.WHITE, shadow=True)

    # Graph title
    add_text_box(slide,
                 Inches(1.0), Inches(1.7), Inches(3), Inches(0.4),
                 "VIEWS",
                 font_name=Fonts.BODY, font_size=Pt(11),
                 font_color=Colors.TEXT_MUTED, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # Y-axis line
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(1.5), Inches(2.2), Inches(0.02), Inches(2.3),
              Colors.TEXT_LIGHT)

    # X-axis line
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(1.5), Inches(4.5), Inches(4.8), Inches(0.02),
              Colors.TEXT_LIGHT)

    # The flatline - dramatic red line at bottom
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(1.5), Inches(4.4), Inches(4.8), Inches(0.04),
              Colors.CORAL)

    # "ZERO" label
    add_text_box(slide,
                 Inches(3.5), Inches(3.8), Inches(2), Inches(0.5),
                 "ZERO",
                 font_name=Fonts.BODY, font_size=Pt(18),
                 font_color=Colors.CORAL, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Quote
    add_text_box(slide,
                 Inches(0.75), Inches(5.3), Inches(6), Inches(0.8),
                 "\"Not fewer views. Zero. Like a switch was flipped.\"",
                 font_name=Fonts.BODY, font_size=Pt(14),
                 font_color=Colors.TEXT_MUTED, italic=True,
                 alignment=PP_ALIGN.LEFT)

    return slide


def build_slide_15_stat_47(prs):
    """
    SLIDE 15: 47+ don't know what to design

    Design: Confusion visualization - scattered question marks
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.CORAL, Inches(0.06))

    # BIG NUMBER - centered this time
    add_text_box(slide,
                 Inches(1), Inches(1.0), Inches(11.333), Inches(1.8),
                 "47+",
                 font_name=Fonts.HEADER, font_size=Pt(96),
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Context
    add_text_box(slide,
                 Inches(1), Inches(2.7), Inches(11.333), Inches(1.0),
                 "said they have \"No idea what to design\"",
                 font_name=Fonts.BODY, font_size=Pt(26),
                 font_color=Colors.TEXT_DARK,
                 alignment=PP_ALIGN.CENTER)

    # Scattered question marks - visual chaos representing confusion
    import random
    random.seed(42)  # Consistent randomness
    positions = [
        (1.5, 4.0), (3.0, 4.5), (4.5, 3.8), (6.0, 4.3), (7.5, 4.0),
        (9.0, 4.4), (10.5, 4.1), (2.0, 5.0), (4.0, 5.2), (6.5, 5.1),
        (8.5, 5.0), (10.0, 5.3), (3.5, 5.8), (5.5, 5.6), (7.8, 5.7)
    ]
    sizes = [Pt(24), Pt(32), Pt(28), Pt(20), Pt(36)]

    for i, (x, y) in enumerate(positions):
        size = sizes[i % len(sizes)]
        alpha_color = Colors.TEXT_LIGHT if i % 2 == 0 else Colors.CORAL_SOFT
        add_text_box(slide,
                     Inches(x), Inches(y), Inches(0.5), Inches(0.6),
                     "?",
                     font_name=Fonts.HEADER, font_size=size,
                     font_color=alpha_color,
                     alignment=PP_ALIGN.CENTER)

    # Empathy card
    empathy_card = add_rounded_rectangle(slide,
                                          Inches(2.5), Inches(6.0),
                                          Inches(8.333), Inches(1.0),
                                          Colors.TEAL_DEEP, shadow=False)

    add_text_box(slide,
                 Inches(2.7), Inches(6.15), Inches(7.933), Inches(0.7),
                 "They're not lazy. They're paralyzed because the rules changed.",
                 font_name=Fonts.BODY, font_size=Pt(16),
                 font_color=Colors.WHITE,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_16_feel_that_way(prs):
    """
    SLIDE 16: Does anyone else feel that way? Type YES

    Design: Engagement prompt - warm, inviting
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.BLUSH)

    # Chat bubble design
    bubble = add_rounded_rectangle(slide,
                                    Inches(2.5), Inches(1.8),
                                    Inches(8.333), Inches(3.0),
                                    Colors.WHITE, shadow=True)

    # Main question
    add_text_box(slide,
                 Inches(2.8), Inches(2.2), Inches(7.733), Inches(1.0),
                 "Does anyone else feel that way?",
                 font_name=Fonts.HEADER, font_size=Pt(32),
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # CTA badge
    yes_badge = add_rounded_rectangle(slide,
                                       Inches(5.0), Inches(3.5),
                                       Inches(3.333), Inches(0.8),
                                       Colors.TEAL_DEEP, shadow=False)

    yes_tf = yes_badge.text_frame
    yes_tf.clear()
    p = yes_tf.paragraphs[0]
    p.text = "Type \"YES\" in chat"
    p.font.name = Fonts.BODY
    p.font.size = Pt(18)
    p.font.color.rgb = Colors.WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    yes_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Subtext
    add_text_box(slide,
                 Inches(2.5), Inches(5.2), Inches(8.333), Inches(0.6),
                 "I'm curious how many of you are seeing the same problems",
                 font_name=Fonts.BODY, font_size=Pt(16),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_17_real_question(prs):
    """
    SLIDE 17: So the real question isn't: "how do I make more listings?"

    Design: Crossing out the wrong question
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.TEAL_DEEP, Inches(0.06))

    # Setup
    add_text_box(slide,
                 Inches(1), Inches(2.0), Inches(11.333), Inches(0.8),
                 "So the real question isn't:",
                 font_name=Fonts.BODY, font_size=Pt(24),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER)

    # The wrong question - crossed out
    wrong_q = add_text_box(slide,
                            Inches(1), Inches(3.2), Inches(11.333), Inches(1.2),
                            "\"How do I make more listings?\"",
                            font_name=Fonts.HEADER, font_size=Pt(40),
                            font_color=Colors.TEXT_LIGHT,
                            alignment=PP_ALIGN.CENTER)

    # Strikethrough line
    strike = add_shape(slide, MSO_SHAPE.RECTANGLE,
                       Inches(3.0), Inches(3.75), Inches(7.333), Inches(0.04),
                       Colors.CORAL)

    # X marks
    add_text_box(slide,
                 Inches(2.0), Inches(3.3), Inches(1), Inches(1),
                 "X",
                 font_name=Fonts.HEADER, font_size=Pt(36),
                 font_color=Colors.CORAL, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide,
                 Inches(10.333), Inches(3.3), Inches(1), Inches(1),
                 "X",
                 font_name=Fonts.HEADER, font_size=Pt(36),
                 font_color=Colors.CORAL, bold=True,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_18_worth_pursuing(prs):
    """
    SLIDE 18: Is Etsy even worth pursuing in 2026?

    Design: THE question - dramatic, full impact
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.TEXT_DARK)

    # The big question - white on dark
    add_text_box(slide,
                 Inches(1), Inches(2.5), Inches(11.333), Inches(2.5),
                 "Is Etsy even worth\npursuing in 2026?",
                 font_name=Fonts.HEADER, font_size=Pt(54),
                 font_color=Colors.WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtle accent
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(5.5), Inches(5.5), Inches(2.333), Inches(0.06),
              Colors.CORAL)

    return slide


def build_slide_19_answer_tonight(prs):
    """
    SLIDE 19: I'm going to answer that question tonight

    Design: Promise, anticipation building
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.TEAL_DEEP, Inches(0.06))
    add_accent_bar(slide, 'bottom', Colors.CORAL, Inches(0.06))

    # Main statement
    add_text_box(slide,
                 Inches(1), Inches(2.0), Inches(11.333), Inches(1.0),
                 "I'm going to answer that question tonight.",
                 font_name=Fonts.HEADER, font_size=Pt(36),
                 font_color=Colors.TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Teaser card
    teaser_card = add_rounded_rectangle(slide,
                                         Inches(3.5), Inches(3.5),
                                         Inches(6.333), Inches(2.0),
                                         Colors.TEAL_DEEP, shadow=True)

    add_text_box(slide,
                 Inches(3.7), Inches(3.8), Inches(5.933), Inches(0.6),
                 "And the answer?",
                 font_name=Fonts.BODY, font_size=Pt(18),
                 font_color=Colors.WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide,
                 Inches(3.7), Inches(4.4), Inches(5.933), Inches(0.8),
                 "It might honestly surprise you.",
                 font_name=Fonts.HEADER, font_size=Pt(28),
                 font_color=Colors.WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    return slide


def build_slide_20_ai_opportunity(prs):
    """
    SLIDE 20: The AI flood actually CREATES an opportunity

    Design: Pattern interrupt - revealing the twist
    """
    slide = add_blank_slide(prs)
    set_background_color(slide, Colors.CREAM)

    add_accent_bar(slide, 'top', Colors.CORAL, Inches(0.08))

    # Opening
    add_text_box(slide,
                 Inches(1), Inches(0.8), Inches(11.333), Inches(0.6),
                 "Because what nobody is talking about is this:",
                 font_name=Fonts.BODY, font_size=Pt(18),
                 font_color=Colors.TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER)

    # Main reveal card
    reveal_card = add_rounded_rectangle(slide,
                                         Inches(1.5), Inches(1.6),
                                         Inches(10.333), Inches(2.5),
                                         Colors.TEAL_DEEP, shadow=True)

    add_text_box(slide,
                 Inches(1.7), Inches(1.9), Inches(9.933), Inches(0.8),
                 "The AI flood actually",
                 font_name=Fonts.BODY, font_size=Pt(22),
                 font_color=Colors.WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide,
                 Inches(1.7), Inches(2.6), Inches(9.933), Inches(1.0),
                 "CREATES an opportunity",
                 font_name=Fonts.HEADER, font_size=Pt(42),
                 font_color=Colors.WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # For who?
    add_text_box(slide,
                 Inches(1), Inches(4.3), Inches(11.333), Inches(0.6),
                 "— for a very specific type of seller...",
                 font_name=Fonts.BODY, font_size=Pt(20),
                 font_color=Colors.TEXT_DARK,
                 alignment=PP_ALIGN.CENTER)

    # BUT warning card
    but_card = add_rounded_rectangle(slide,
                                      Inches(4.5), Inches(5.2),
                                      Inches(4.333), Inches(1.0),
                                      Colors.CORAL, shadow=False)

    add_text_box(slide,
                 Inches(4.5), Inches(5.4), Inches(4.333), Inches(0.6),
                 "BUT",
                 font_name=Fonts.HEADER, font_size=Pt(28),
                 font_color=Colors.WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Condition
    add_text_box(slide,
                 Inches(1), Inches(6.4), Inches(11.333), Inches(0.6),
                 "Only if you understand what's really happening...",
                 font_name=Fonts.BODY, font_size=Pt(18),
                 font_color=Colors.TEXT_MUTED, italic=True,
                 alignment=PP_ALIGN.CENTER)

    return slide


# =============================================================================
# MAIN BUILD FUNCTION
# =============================================================================

def build_presentation():
    """Build the complete presentation"""
    print("Creating Bailey Vann - The 2026 Etsy Reset presentation...")

    prs = create_presentation()

    # Build slides 1-20
    print("Building Slide 1: Title...")
    build_slide_1_title(prs)

    print("Building Slide 2: Before We Begin...")
    build_slide_2_before_we_begin(prs)

    print("Building Slide 3: Get Ready Chat...")
    build_slide_3_get_ready_chat(prs)

    print("Building Slide 4: Quiz A/B...")
    build_slide_4_quiz_ab(prs)

    print("Building Slide 5: Type A or B...")
    build_slide_5_type_ab(prs)

    print("Building Slide 6: The Answer Is...")
    build_slide_6_answer_is(prs)

    print("Building Slide 7: Both AI...")
    build_slide_7_both_ai(prs)

    print("Building Slide 8: AI Nowadays...")
    build_slide_8_ai_nowadays(prs)

    print("Building Slide 9: Sink In...")
    build_slide_9_sink_in(prs)

    print("Building Slide 10: Uncomfortable...")
    build_slide_10_uncomfortable(prs)

    print("Building Slide 11: What Happens...")
    build_slide_11_what_happens(prs)

    print("Building Slide 12: Survey Intro...")
    build_slide_12_survey_intro(prs)

    print("Building Slide 13: Stat 75+...")
    build_slide_13_stat_75(prs)

    print("Building Slide 14: Stat 58+...")
    build_slide_14_stat_58(prs)

    print("Building Slide 15: Stat 47+...")
    build_slide_15_stat_47(prs)

    print("Building Slide 16: Feel That Way...")
    build_slide_16_feel_that_way(prs)

    print("Building Slide 17: Real Question...")
    build_slide_17_real_question(prs)

    print("Building Slide 18: Worth Pursuing...")
    build_slide_18_worth_pursuing(prs)

    print("Building Slide 19: Answer Tonight...")
    build_slide_19_answer_tonight(prs)

    print("Building Slide 20: AI Opportunity...")
    build_slide_20_ai_opportunity(prs)

    # Save
    output_path = "/home/user/webby-slides-bailey/Bailey_Etsy_Reset_Redesign_v1.pptx"
    prs.save(output_path)
    print(f"\nSaved to: {output_path}")
    print("Done! 20 slides created.")

    return output_path


if __name__ == "__main__":
    build_presentation()
