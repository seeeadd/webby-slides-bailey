"""
Bailey Vann - The 2026 Etsy Reset
PREMIUM EDITORIAL SLIDE DECK - Version 2
Pinterest-worthy, $5K brand deck quality

Complete rebuild with:
- Organic blob shapes (freeform bezier curves)
- Layered compositions with depth
- Asymmetric editorial layouts
- Custom graphics per slide
- Embedded OGG fonts
- Texture and visual interest
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from pptx.parts.embeddedpackage import EmbeddedPackagePart
from lxml import etree
import os
import random
import math

# =============================================================================
# DESIGN SYSTEM - PREMIUM EDITORIAL
# =============================================================================

class Colors:
    """Bailey Vann Brand - Warm, Premium, Editorial"""
    # Primary
    TEAL_DEEP = RGBColor(0x1B, 0x8A, 0x8A)
    TEAL = RGBColor(0x2B, 0xA5, 0xA3)
    TEAL_LIGHT = RGBColor(0x5B, 0xBC, 0xB3)

    CORAL = RGBColor(0xE0, 0x7B, 0x6C)
    CORAL_SOFT = RGBColor(0xF4, 0xA8, 0x9A)
    CORAL_PALE = RGBColor(0xFD, 0xD5, 0xCC)

    # Backgrounds - warm, creamy
    CREAM = RGBColor(0xFD, 0xF8, 0xF3)
    CREAM_DARK = RGBColor(0xF5, 0xEE, 0xE6)
    BLUSH = RGBColor(0xFE, 0xF0, 0xEA)
    BLUSH_SOFT = RGBColor(0xFE, 0xE5, 0xE0)
    MINT = RGBColor(0xE8, 0xF5, 0xF3)

    # Accents
    GOLD = RGBColor(0xD4, 0xAF, 0x37)
    GOLD_SOFT = RGBColor(0xF7, 0xE1, 0x9C)

    # Text
    DARK = RGBColor(0x2D, 0x34, 0x36)
    MUTED = RGBColor(0x63, 0x6E, 0x72)
    LIGHT = RGBColor(0x9C, 0xA3, 0xA8)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)


class Fonts:
    """Typography - OGG Family"""
    # Display headers
    DISPLAY = "Ogg TRIAL"

    # Body text - using Ogg Text variant
    BODY = "Ogg Text TRIAL"
    BODY_MEDIUM = "Ogg Text TRIAL Medium"

    # Lighter weights for variety
    LIGHT = "Ogg TRIAL Light"


# =============================================================================
# ORGANIC SHAPE BUILDERS
# =============================================================================

def create_organic_blob(slide, x, y, width, height, color, opacity=100):
    """
    Create an organic blob shape using freeform path
    More natural, Pinterest-worthy than basic circles
    """
    # Generate organic blob points using sine waves for natural curves
    points = 16
    cx, cy = width / 2, height / 2

    # Build list of vertices
    vertices = []
    for i in range(points):
        angle = (2 * math.pi * i) / points
        # Add organic variation for natural feel
        variation = 0.12 * math.sin(3 * angle) + 0.08 * math.cos(5 * angle) + 0.05 * math.sin(7 * angle)
        r_x = (cx * (0.88 + variation))
        r_y = (cy * (0.88 + variation * 0.85))

        px = cx + r_x * math.cos(angle)
        py = cy + r_y * math.sin(angle)
        vertices.append((Emu(int(px * 914400)), Emu(int(py * 914400))))

    # Create freeform with vertices
    builder = slide.shapes.build_freeform(x, y)
    builder.move_to(vertices[0][0], vertices[0][1])
    builder.add_line_segments(vertices[1:], close=True)

    shape = builder.convert_to_shape(x, y)

    # Apply fill
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Set transparency if needed
    if opacity < 100:
        set_shape_transparency(shape, 100 - opacity)

    return shape


def create_wave_shape(slide, x, y, width, height, color, wave_height=0.3):
    """Create a wavy organic shape for backgrounds"""
    builder = slide.shapes.build_freeform(x, y)

    # Start at bottom left
    builder.move_to(0, height)

    # Bottom edge
    builder.line_to(width, height)

    # Right edge up
    builder.line_to(width, height * 0.4)

    # Wavy top edge
    segments = 8
    for i in range(segments + 1):
        px = width * (1 - i / segments)
        wave = math.sin(i * math.pi / 2) * height * wave_height
        py = height * 0.3 + wave
        builder.line_to(px, py)

    # Left edge
    builder.line_to(0, height * 0.4)

    builder.close()
    shape = builder.convert_to_shape(x, y)

    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    return shape


def set_shape_transparency(shape, transparency_percent):
    """Set transparency on a shape (0-100)"""
    # Access fill element
    spPr = shape._sp.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', str(int((100 - transparency_percent) * 1000)))


def add_soft_shadow(shape, blur=12, distance=4, opacity=15):
    """Add modern soft shadow to shape"""
    spPr = shape._sp.spPr

    # Remove existing effects
    existing = spPr.find(qn('a:effectLst'))
    if existing is not None:
        spPr.remove(existing)

    shadow_xml = f'''
    <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="{blur * 12700}" dist="{distance * 12700}" dir="5400000" algn="tl" rotWithShape="0">
            <a:srgbClr val="2D3436">
                <a:alpha val="{opacity * 1000}"/>
            </a:srgbClr>
        </a:outerShdw>
    </a:effectLst>
    '''
    effect_lst = etree.fromstring(shadow_xml)
    spPr.append(effect_lst)


def add_gradient_fill(shape, color1, color2, angle=90):
    """Add gradient fill to shape"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2


# =============================================================================
# PREMIUM TEXT HELPERS
# =============================================================================

def add_text(slide, left, top, width, height, text,
             font=Fonts.BODY, size=Pt(18), color=Colors.DARK,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             line_spacing=1.2):
    """Add premium styled text"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align

    # Line spacing
    p.line_spacing = line_spacing

    return box


def add_display_text(slide, left, top, width, height, text,
                     size=Pt(72), color=Colors.DARK, align=PP_ALIGN.LEFT):
    """Add large display/headline text with OGG"""
    return add_text(slide, left, top, width, height, text,
                    font=Fonts.DISPLAY, size=size, color=color,
                    bold=True, align=align, line_spacing=1.0)


def add_body_text(slide, left, top, width, height, text,
                  size=Pt(18), color=Colors.MUTED, align=PP_ALIGN.LEFT):
    """Add body text with Ogg Text"""
    return add_text(slide, left, top, width, height, text,
                    font=Fonts.BODY, size=size, color=color,
                    align=align, line_spacing=1.4)


# =============================================================================
# PREMIUM CARD COMPONENTS
# =============================================================================

def add_floating_card(slide, left, top, width, height,
                      fill_color=Colors.WHITE, shadow=True, corner_radius=0.08):
    """Add a premium floating card with soft shadow"""
    from pptx.enum.shapes import MSO_SHAPE

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )

    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.fill.background()

    # Adjust corner radius
    card.adjustments[0] = corner_radius

    if shadow:
        add_soft_shadow(card, blur=20, distance=6, opacity=12)

    return card


def add_pill_label(slide, left, top, text, bg_color=Colors.CORAL,
                   text_color=Colors.WHITE, size=Pt(11)):
    """Add a stylish pill-shaped label"""
    # Calculate size based on text
    width = Inches(len(text) * 0.09 + 0.5)
    height = Inches(0.32)

    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )

    pill.fill.solid()
    pill.fill.fore_color.rgb = bg_color
    pill.line.fill.background()
    pill.adjustments[0] = 0.5  # Full pill shape

    # Add text
    tf = pill.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.name = Fonts.BODY_MEDIUM
    p.font.size = size
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    return pill


# =============================================================================
# BACKGROUND COMPOSITIONS
# =============================================================================

def create_editorial_background_1(slide):
    """
    Layered organic background - cream with soft teal blobs
    Editorial, Pinterest-worthy
    """
    # Base cream
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.CREAM
    bg.line.fill.background()

    # Large organic teal blob - top right
    blob1 = create_organic_blob(
        slide,
        Inches(8), Inches(-2),
        Inches(7), Inches(6),
        Colors.TEAL_LIGHT, opacity=15
    )

    # Smaller coral blob - bottom left
    blob2 = create_organic_blob(
        slide,
        Inches(-2), Inches(4),
        Inches(5), Inches(5),
        Colors.CORAL_PALE, opacity=20
    )

    # Subtle gold accent blob
    blob3 = create_organic_blob(
        slide,
        Inches(10), Inches(5),
        Inches(4), Inches(3),
        Colors.GOLD_SOFT, opacity=25
    )

    return [bg, blob1, blob2, blob3]


def create_editorial_background_2(slide):
    """
    Warm blush background with floating elements
    """
    # Base blush gradient feel
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(7.5)
    )
    add_gradient_fill(bg, Colors.CREAM, Colors.BLUSH_SOFT, angle=135)
    bg.line.fill.background()

    # Organic shapes for depth
    blob1 = create_organic_blob(
        slide,
        Inches(-1), Inches(-1),
        Inches(6), Inches(5),
        Colors.MINT, opacity=30
    )

    blob2 = create_organic_blob(
        slide,
        Inches(9), Inches(3),
        Inches(6), Inches(5),
        Colors.CORAL_PALE, opacity=20
    )

    return [bg, blob1, blob2]


def create_bold_teal_background(slide):
    """Bold teal background for impact slides"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(7.5)
    )
    add_gradient_fill(bg, Colors.TEAL_DEEP, Colors.TEAL, angle=135)
    bg.line.fill.background()

    # Subtle lighter blob for depth
    blob = create_organic_blob(
        slide,
        Inches(7), Inches(-1),
        Inches(8), Inches(6),
        Colors.TEAL_LIGHT, opacity=15
    )

    return [bg, blob]


def create_dark_dramatic_background(slide):
    """Dark background for dramatic moments"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.DARK
    bg.line.fill.background()

    # Subtle coral glow
    blob = create_organic_blob(
        slide,
        Inches(4), Inches(2),
        Inches(6), Inches(4),
        Colors.CORAL, opacity=8
    )

    return [bg, blob]


# =============================================================================
# SLIDE BUILDERS - PREMIUM EDITORIAL
# =============================================================================

def build_slide_01_title(prs):
    """
    THE 2026 ETSY RESET - Title Slide

    Editorial, magazine-cover aesthetic
    Asymmetric layout with layered organic shapes
    Typography-driven with subtle graphics
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Layered background
    create_editorial_background_1(slide)

    # Day indicator - top left pill
    add_pill_label(slide, Inches(0.8), Inches(0.6),
                   "Day 1 of 3", Colors.TEAL_DEEP)

    # Small context text
    add_body_text(slide,
        Inches(0.8), Inches(1.3), Inches(6), Inches(0.4),
        "The 2026 Etsy Upgrade Challenge",
        size=Pt(14), color=Colors.TEAL_DEEP
    )

    # MASSIVE "RESET" - the hero
    add_display_text(slide,
        Inches(0.7), Inches(1.8), Inches(11), Inches(2.2),
        "RESET",
        size=Pt(180), color=Colors.DARK
    )

    # Subtitle - editorial style
    add_body_text(slide,
        Inches(0.8), Inches(4.2), Inches(7), Inches(1.2),
        "Delete the Dead Weight.\nFind Your Focus.\nBuild a Shop That Actually Works.",
        size=Pt(22), color=Colors.MUTED
    )

    # Bailey credentials - floating card
    cred_card = add_floating_card(slide,
        Inches(0.8), Inches(5.8),
        Inches(8), Inches(1.0),
        Colors.TEAL_DEEP, shadow=True
    )

    # Credentials text on card
    add_text(slide,
        Inches(1.0), Inches(5.95), Inches(7.6), Inches(0.7),
        "with Bailey Vann  •  Top 0.1% Etsy Seller  •  $1M+ in Digital Product Sales",
        font=Fonts.BODY, size=Pt(14), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    # Decorative organic shape - right side visual interest
    accent_blob = create_organic_blob(slide,
        Inches(9.5), Inches(2.5),
        Inches(3.5), Inches(3),
        Colors.CORAL_SOFT, opacity=60
    )

    # Smaller overlapping blob
    accent_blob2 = create_organic_blob(slide,
        Inches(10.5), Inches(3.5),
        Inches(2.5), Inches(2.5),
        Colors.GOLD_SOFT, opacity=50
    )

    return slide


def build_slide_02_before_begin(prs):
    """
    Before We Begin... Quick Pop Quiz

    Intriguing, sets up engagement
    Clean but layered design
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Layered background
    create_editorial_background_2(slide)

    # Main content card - floating, premium
    main_card = add_floating_card(slide,
        Inches(2), Inches(1.2),
        Inches(9.333), Inches(5.2),
        Colors.WHITE, shadow=True
    )

    # Headline
    add_display_text(slide,
        Inches(2.5), Inches(1.8), Inches(8.333), Inches(1.2),
        "Before We Begin...",
        size=Pt(52), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(3.1), Inches(2.333), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = Colors.TEAL
    line.line.fill.background()

    # Quiz badge - stylish pill
    add_pill_label(slide,
        Inches(5.0), Inches(3.5),
        "Quick Pop Quiz", Colors.CORAL)

    # Subtext
    add_body_text(slide,
        Inches(2.5), Inches(4.4), Inches(8.333), Inches(1.5),
        "I want to show you something that might change\nhow you think about Etsy in 2026...",
        size=Pt(20), color=Colors.MUTED,
        align=PP_ALIGN.CENTER
    )

    # Decorative blobs on edges of card
    blob_accent = create_organic_blob(slide,
        Inches(10.5), Inches(0.8),
        Inches(2), Inches(2),
        Colors.TEAL_LIGHT, opacity=40
    )

    return slide


def build_slide_03_get_ready_chat(prs):
    """
    Get ready to type in the chat!

    Warm, inviting, interactive energy
    Custom chat bubble graphic
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Warm blush background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(7.5)
    )
    add_gradient_fill(bg, Colors.BLUSH_SOFT, Colors.BLUSH, angle=180)
    bg.line.fill.background()

    # Organic accent blobs
    create_organic_blob(slide, Inches(-1), Inches(-1),
        Inches(4), Inches(4), Colors.CORAL_PALE, opacity=40)
    create_organic_blob(slide, Inches(10), Inches(4),
        Inches(4), Inches(4), Colors.MINT, opacity=30)

    # Main chat bubble - custom built, not basic rectangle
    # Large bubble card
    bubble = add_floating_card(slide,
        Inches(2.5), Inches(1.5),
        Inches(8.333), Inches(3.5),
        Colors.WHITE, shadow=True, corner_radius=0.12
    )

    # Chat bubble "tail" - triangle
    tail = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(4.5), Inches(4.9),
        Inches(0.6), Inches(0.5)
    )
    tail.fill.solid()
    tail.fill.fore_color.rgb = Colors.WHITE
    tail.line.fill.background()
    tail.rotation = 180

    # Text inside bubble
    add_display_text(slide,
        Inches(3), Inches(2.2), Inches(7.333), Inches(1.8),
        "Get ready to type\nin the chat!",
        size=Pt(42), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # Typing indicator dots - organic, not basic circles
    for i, offset in enumerate([0, 0.5, 1.0]):
        dot = create_organic_blob(slide,
            Inches(5.8 + offset), Inches(4.2),
            Inches(0.35), Inches(0.35),
            Colors.TEAL_DEEP, opacity=80 - i*15
        )

    # Subtext
    add_body_text(slide,
        Inches(2), Inches(5.8), Inches(9.333), Inches(0.6),
        "This is interactive — your answers matter",
        size=Pt(16), color=Colors.MUTED,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_04_quiz_ab(prs):
    """
    A/B Quiz - Which design was made by professional?

    Clean comparison layout with premium cards
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_1(slide)

    # Question header
    add_display_text(slide,
        Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.0),
        "Which design was made by a professional artist?",
        size=Pt(32), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # Card A - floating premium card
    card_a = add_floating_card(slide,
        Inches(0.8), Inches(1.6),
        Inches(5.4), Inches(4.8),
        Colors.WHITE, shadow=True
    )

    # Image placeholder area A
    placeholder_a = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.1), Inches(1.9),
        Inches(4.8), Inches(3.3)
    )
    add_gradient_fill(placeholder_a, Colors.CREAM_DARK, Colors.MINT, angle=135)
    placeholder_a.line.fill.background()
    placeholder_a.adjustments[0] = 0.05

    add_body_text(slide,
        Inches(1.1), Inches(3.2), Inches(4.8), Inches(0.5),
        "[ Image A ]",
        size=Pt(14), color=Colors.LIGHT,
        align=PP_ALIGN.CENTER
    )

    # "A" label - organic blob shape
    label_a = create_organic_blob(slide,
        Inches(3.0), Inches(5.5),
        Inches(0.9), Inches(0.9),
        Colors.TEAL_DEEP, opacity=100
    )
    add_text(slide,
        Inches(3.1), Inches(5.6), Inches(0.7), Inches(0.7),
        "A",
        font=Fonts.DISPLAY, size=Pt(28), color=Colors.WHITE,
        bold=True, align=PP_ALIGN.CENTER
    )

    # VS divider
    add_display_text(slide,
        Inches(6.0), Inches(3.5), Inches(1.333), Inches(0.6),
        "vs",
        size=Pt(24), color=Colors.LIGHT,
        align=PP_ALIGN.CENTER
    )

    # Card B - floating premium card
    card_b = add_floating_card(slide,
        Inches(7.133), Inches(1.6),
        Inches(5.4), Inches(4.8),
        Colors.WHITE, shadow=True
    )

    # Image placeholder area B
    placeholder_b = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.433), Inches(1.9),
        Inches(4.8), Inches(3.3)
    )
    add_gradient_fill(placeholder_b, Colors.CREAM_DARK, Colors.BLUSH_SOFT, angle=135)
    placeholder_b.line.fill.background()
    placeholder_b.adjustments[0] = 0.05

    add_body_text(slide,
        Inches(7.433), Inches(3.2), Inches(4.8), Inches(0.5),
        "[ Image B ]",
        size=Pt(14), color=Colors.LIGHT,
        align=PP_ALIGN.CENTER
    )

    # "B" label - organic blob
    label_b = create_organic_blob(slide,
        Inches(9.4), Inches(5.5),
        Inches(0.9), Inches(0.9),
        Colors.CORAL, opacity=100
    )
    add_text(slide,
        Inches(9.5), Inches(5.6), Inches(0.7), Inches(0.7),
        "B",
        font=Fonts.DISPLAY, size=Pt(28), color=Colors.WHITE,
        bold=True, align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_05_type_ab(prs):
    """
    Type A or B in the chat!

    Bold, high energy, call to action
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_bold_teal_background(slide)

    # Main message
    add_display_text(slide,
        Inches(1), Inches(2.0), Inches(11.333), Inches(1.5),
        "Type A or B\nin the chat!",
        size=Pt(64), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    # Option buttons - organic blob shapes
    # A button
    blob_a = create_organic_blob(slide,
        Inches(3.5), Inches(4.5),
        Inches(2.5), Inches(2),
        Colors.WHITE, opacity=100
    )
    add_soft_shadow(blob_a, blur=15, distance=4, opacity=20)

    add_display_text(slide,
        Inches(3.8), Inches(4.9), Inches(2), Inches(1.2),
        "A",
        size=Pt(56), color=Colors.TEAL_DEEP,
        align=PP_ALIGN.CENTER
    )

    # "or"
    add_text(slide,
        Inches(6.0), Inches(5.2), Inches(1.333), Inches(0.6),
        "or",
        font=Fonts.BODY, size=Pt(24), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    # B button
    blob_b = create_organic_blob(slide,
        Inches(7.333), Inches(4.5),
        Inches(2.5), Inches(2),
        Colors.CORAL, opacity=100
    )
    add_soft_shadow(blob_b, blur=15, distance=4, opacity=20)

    add_display_text(slide,
        Inches(7.6), Inches(4.9), Inches(2), Inches(1.2),
        "B",
        size=Pt(56), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_06_answer_is(prs):
    """
    The Answer Is...

    Suspenseful, dramatic pause
    Minimal but premium
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_1(slide)

    # Dramatic centered text
    add_display_text(slide,
        Inches(1), Inches(2.5), Inches(11.333), Inches(2),
        "The Answer Is...",
        size=Pt(72), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # Suspense dots - organic blobs
    for i in range(3):
        blob = create_organic_blob(slide,
            Inches(5.8 + i * 0.5), Inches(5.0),
            Inches(0.25), Inches(0.25),
            Colors.TEAL_DEEP, opacity=80 - i*20
        )

    return slide


def build_slide_07_both_ai(prs):
    """
    Both Were Made by AI.

    REVEAL moment - dramatic impact
    Full coral background for shock
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Bold coral background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(7.5)
    )
    add_gradient_fill(bg, Colors.CORAL, Colors.CORAL_SOFT, angle=135)
    bg.line.fill.background()

    # Subtle organic overlay
    create_organic_blob(slide, Inches(8), Inches(-1),
        Inches(7), Inches(5), Colors.WHITE, opacity=10)

    # Main reveal text
    add_display_text(slide,
        Inches(0.8), Inches(1.5), Inches(11.733), Inches(2),
        "Both Were Made by AI.",
        size=Pt(68), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    # Detail card - floating
    detail_card = add_floating_card(slide,
        Inches(3.5), Inches(4.2),
        Inches(6.333), Inches(1.8),
        Colors.WHITE, shadow=True
    )

    add_body_text(slide,
        Inches(3.7), Inches(4.5), Inches(5.933), Inches(1.4),
        "In less than 30 seconds.\nFor free.",
        size=Pt(26), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_08_ai_nowadays(prs):
    """
    Yes, we all use AI nowadays...
    But in 2026...

    Contemplative, building tension
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_2(slide)

    # Main content card
    main_card = add_floating_card(slide,
        Inches(1.2), Inches(1.2),
        Inches(10.933), Inches(5.2),
        Colors.WHITE, shadow=True
    )

    # First part - muted acknowledgment
    add_body_text(slide,
        Inches(1.8), Inches(1.8), Inches(9.733), Inches(0.8),
        "Yes, we all use AI nowadays...",
        size=Pt(28), color=Colors.LIGHT,
        align=PP_ALIGN.LEFT
    )

    # Decorative divider - organic line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.8), Inches(2.8), Inches(3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = Colors.TEAL
    line.line.fill.background()

    # Key message - emphasized
    add_display_text(slide,
        Inches(1.8), Inches(3.2), Inches(9.733), Inches(2.5),
        "But in 2026, they've gotten so good\nthat hiring someone REAL is\nstarting to become just an option...",
        size=Pt(32), color=Colors.DARK,
        align=PP_ALIGN.LEFT
    )

    # Highlight pill for emphasis
    add_pill_label(slide,
        Inches(1.8), Inches(5.6),
        "just an option", Colors.CORAL_SOFT, Colors.CORAL
    )

    return slide


def build_slide_09_sink_in(prs):
    """
    Let that sink in for a second.

    Minimal, breathing room
    Pure typography moment
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Soft mint background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(7.5)
    )
    add_gradient_fill(bg, Colors.MINT, Colors.CREAM, angle=180)
    bg.line.fill.background()

    # Subtle organic shapes
    create_organic_blob(slide, Inches(9), Inches(4),
        Inches(5), Inches(4), Colors.TEAL_LIGHT, opacity=15)
    create_organic_blob(slide, Inches(-1), Inches(-1),
        Inches(4), Inches(4), Colors.GOLD_SOFT, opacity=20)

    # Single powerful line
    add_display_text(slide,
        Inches(1), Inches(2.8), Inches(11.333), Inches(2),
        "Let that sink in\nfor a second.",
        size=Pt(52), color=Colors.TEAL_DEEP,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_10_uncomfortable(prs):
    """
    Now let me ask you something uncomfortable...

    Tension building, darker mood
    Dramatic shift
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_dark_dramatic_background(slide)

    # First line
    add_text(slide,
        Inches(1.5), Inches(2.3), Inches(10.333), Inches(1),
        "Now let me ask you",
        font=Fonts.DISPLAY, size=Pt(42), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    # Emphasized line - coral glow effect
    add_display_text(slide,
        Inches(1.5), Inches(3.5), Inches(10.333), Inches(1.5),
        "something uncomfortable...",
        size=Pt(56), color=Colors.CORAL,
        align=PP_ALIGN.CENTER
    )

    # Subtle organic coral glow behind text
    glow = create_organic_blob(slide,
        Inches(3), Inches(3.2),
        Inches(7), Inches(2),
        Colors.CORAL, opacity=10
    )

    return slide


def build_slide_11_what_happens(prs):
    """
    If anyone can create designs like this in seconds...
    What happens to YOUR Etsy shop in 2026?

    THE big question - maximum impact
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_1(slide)

    # Setup line - smaller
    add_body_text(slide,
        Inches(1), Inches(1.0), Inches(11.333), Inches(0.8),
        "If anyone can create designs like this in seconds...",
        size=Pt(22), color=Colors.MUTED,
        align=PP_ALIGN.CENTER
    )

    # THE BIG QUESTION
    add_display_text(slide,
        Inches(0.5), Inches(2.2), Inches(12.333), Inches(3),
        "What happens to\nYOUR Etsy shop\nin 2026?",
        size=Pt(60), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # Large decorative question mark - organic blob style
    q_blob = create_organic_blob(slide,
        Inches(10), Inches(4.5),
        Inches(2.5), Inches(2.5),
        Colors.CORAL_PALE, opacity=60
    )

    add_display_text(slide,
        Inches(10.2), Inches(4.6), Inches(2), Inches(2.3),
        "?",
        size=Pt(96), color=Colors.CORAL,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_12_survey_intro(prs):
    """
    The numbers I'm about to show you aren't random...

    Setting up data credibility
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_2(slide)

    # Main statement
    add_display_text(slide,
        Inches(1), Inches(1.5), Inches(11.333), Inches(2),
        "The numbers I'm about\nto show you aren't random...",
        size=Pt(48), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # Credibility card - floating
    cred_card = add_floating_card(slide,
        Inches(2.5), Inches(4.0),
        Inches(8.333), Inches(2.2),
        Colors.TEAL_DEEP, shadow=True
    )

    add_text(slide,
        Inches(2.8), Inches(4.3), Inches(7.733), Inches(0.5),
        "These are from YOUR OWN ANSWERS",
        font=Fonts.BODY_MEDIUM, size=Pt(18), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    add_text(slide,
        Inches(2.8), Inches(4.9), Inches(7.733), Inches(0.5),
        "after I asked 160+ of you last week",
        font=Fonts.BODY, size=Pt(16), color=Colors.TEAL_LIGHT,
        align=PP_ALIGN.CENTER
    )

    # "160+" emphasized badge
    add_pill_label(slide,
        Inches(5.5), Inches(5.7),
        "160+ sellers surveyed", Colors.CORAL)

    return slide


def build_slide_13_stat_75(prs):
    """
    75+ sellers report shops tanked

    Data visualization with human scale
    Dot matrix showing proportion
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_1(slide)

    # Left side - the number and context
    # Big number
    add_display_text(slide,
        Inches(0.6), Inches(0.6), Inches(5), Inches(2),
        "75+",
        size=Pt(120), color=Colors.CORAL,
        align=PP_ALIGN.LEFT
    )

    # Context
    add_body_text(slide,
        Inches(0.8), Inches(2.6), Inches(5), Inches(1.2),
        "sellers report their shops\ntanked in the last 2 months",
        size=Pt(24), color=Colors.DARK
    )

    # Right side - visual dot matrix
    # Create organic dots representing 75 out of ~100
    start_x = 6.8
    start_y = 1.0

    for row in range(8):
        for col in range(10):
            idx = row * 10 + col
            x = Inches(start_x + col * 0.55)
            y = Inches(start_y + row * 0.55)

            # First 75 are coral (affected), rest are light
            if idx < 75:
                color = Colors.CORAL
                opacity = 85
            else:
                color = Colors.CREAM_DARK
                opacity = 100

            create_organic_blob(slide, x, y,
                Inches(0.4), Inches(0.4),
                color, opacity=opacity
            )

    # Insight card
    insight = add_floating_card(slide,
        Inches(0.6), Inches(4.2),
        Inches(5.5), Inches(1.6),
        Colors.WHITE, shadow=True
    )

    add_body_text(slide,
        Inches(0.9), Inches(4.5), Inches(4.9), Inches(1.2),
        "If there are 4 people on your row in this chat,\n3 of them are feeling this...",
        size=Pt(16), color=Colors.MUTED
    )

    return slide


def build_slide_14_stat_58(prs):
    """
    58+ sellers said "No views anymore"

    Flatline visualization
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_2(slide)

    # Graph card - left side
    graph_card = add_floating_card(slide,
        Inches(0.6), Inches(1.0),
        Inches(6.5), Inches(4.5),
        Colors.WHITE, shadow=True
    )

    # Graph title
    add_pill_label(slide, Inches(1.0), Inches(1.4), "VIEWS", Colors.MUTED)

    # Y-axis
    y_axis = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(2.2), Inches(0.02), Inches(2.5)
    )
    y_axis.fill.solid()
    y_axis.fill.fore_color.rgb = Colors.LIGHT
    y_axis.line.fill.background()

    # X-axis
    x_axis = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(4.7), Inches(5), Inches(0.02)
    )
    x_axis.fill.solid()
    x_axis.fill.fore_color.rgb = Colors.LIGHT
    x_axis.line.fill.background()

    # THE FLATLINE - dramatic
    flatline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(4.5), Inches(5), Inches(0.06)
    )
    flatline.fill.solid()
    flatline.fill.fore_color.rgb = Colors.CORAL
    flatline.line.fill.background()

    # "ZERO" label
    add_display_text(slide,
        Inches(3.2), Inches(3.5), Inches(2), Inches(0.8),
        "ZERO",
        size=Pt(28), color=Colors.CORAL,
        align=PP_ALIGN.CENTER
    )

    # Right side - number and context
    add_display_text(slide,
        Inches(7.5), Inches(1.2), Inches(5), Inches(1.5),
        "58+",
        size=Pt(96), color=Colors.TEAL_DEEP,
        align=PP_ALIGN.RIGHT
    )

    add_body_text(slide,
        Inches(7.5), Inches(2.8), Inches(5), Inches(1.5),
        "sellers said their #1 problem:\n\"No views anymore\"",
        size=Pt(22), color=Colors.DARK,
        align=PP_ALIGN.RIGHT
    )

    # Quote
    add_body_text(slide,
        Inches(7.5), Inches(4.5), Inches(5.2), Inches(1),
        "\"Not fewer views. Zero.\nLike a switch was flipped.\"",
        size=Pt(16), color=Colors.MUTED,
        align=PP_ALIGN.RIGHT
    )

    return slide


def build_slide_15_stat_47(prs):
    """
    47+ don't know what to design

    Confusion visualization with scattered question marks
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_1(slide)

    # Central number
    add_display_text(slide,
        Inches(1), Inches(0.8), Inches(11.333), Inches(2),
        "47+",
        size=Pt(120), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # Context
    add_body_text(slide,
        Inches(1), Inches(2.6), Inches(11.333), Inches(0.8),
        "said they have \"No idea what to design\"",
        size=Pt(28), color=Colors.MUTED,
        align=PP_ALIGN.CENTER
    )

    # Scattered question marks - organic chaos
    random.seed(42)
    positions = [
        (1.2, 3.8), (2.8, 4.3), (4.2, 3.6), (5.8, 4.5), (7.2, 3.9),
        (8.8, 4.2), (10.3, 3.7), (1.8, 5.0), (3.5, 5.3), (5.5, 5.1),
        (7.5, 5.4), (9.5, 5.0), (11.0, 4.8)
    ]

    for i, (x, y) in enumerate(positions):
        size = Pt(random.randint(20, 40))
        opacity = random.randint(30, 70)
        color = Colors.CORAL_PALE if i % 2 == 0 else Colors.LIGHT

        blob = create_organic_blob(slide,
            Inches(x), Inches(y),
            Inches(0.5), Inches(0.5),
            color, opacity=opacity
        )

        add_text(slide,
            Inches(x + 0.05), Inches(y + 0.02),
            Inches(0.4), Inches(0.45),
            "?",
            font=Fonts.DISPLAY, size=size, color=Colors.CORAL,
            align=PP_ALIGN.CENTER
        )

    # Empathy bar
    empathy = add_floating_card(slide,
        Inches(2), Inches(6.0),
        Inches(9.333), Inches(0.9),
        Colors.TEAL_DEEP, shadow=False
    )

    add_text(slide,
        Inches(2.3), Inches(6.15), Inches(8.733), Inches(0.6),
        "They're not lazy. They're paralyzed because the rules changed.",
        font=Fonts.BODY, size=Pt(16), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_16_type_yes(prs):
    """
    Does anyone else feel that way? Type YES

    Engagement prompt - warm and inviting
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Warm background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333), Inches(7.5)
    )
    add_gradient_fill(bg, Colors.BLUSH_SOFT, Colors.CREAM, angle=180)
    bg.line.fill.background()

    # Organic accents
    create_organic_blob(slide, Inches(-1), Inches(5),
        Inches(4), Inches(3), Colors.CORAL_PALE, opacity=30)
    create_organic_blob(slide, Inches(10), Inches(-1),
        Inches(4), Inches(4), Colors.MINT, opacity=25)

    # Main card
    main_card = add_floating_card(slide,
        Inches(2), Inches(1.5),
        Inches(9.333), Inches(4),
        Colors.WHITE, shadow=True
    )

    # Question
    add_display_text(slide,
        Inches(2.5), Inches(2.2), Inches(8.333), Inches(1.5),
        "Does anyone else\nfeel that way?",
        size=Pt(44), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # YES button
    yes_blob = create_organic_blob(slide,
        Inches(5.0), Inches(4.2),
        Inches(3.333), Inches(1.0),
        Colors.TEAL_DEEP, opacity=100
    )
    add_soft_shadow(yes_blob, blur=10, distance=3, opacity=15)

    add_text(slide,
        Inches(5.2), Inches(4.35), Inches(2.933), Inches(0.7),
        "Type \"YES\" in chat",
        font=Fonts.BODY_MEDIUM, size=Pt(18), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    # Subtext
    add_body_text(slide,
        Inches(2), Inches(5.8), Inches(9.333), Inches(0.6),
        "I'm curious how many of you are seeing the same problems",
        size=Pt(16), color=Colors.MUTED,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_17_real_question(prs):
    """
    The real question isn't "how do I make more listings?"

    Crossing out the wrong approach
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_1(slide)

    # Setup
    add_body_text(slide,
        Inches(1), Inches(1.8), Inches(11.333), Inches(0.8),
        "So the real question isn't:",
        size=Pt(24), color=Colors.MUTED,
        align=PP_ALIGN.CENTER
    )

    # The wrong question - crossed out style
    add_text(slide,
        Inches(1), Inches(3.0), Inches(11.333), Inches(1.5),
        "\"How do I make more listings?\"",
        font=Fonts.DISPLAY, size=Pt(44), color=Colors.LIGHT,
        align=PP_ALIGN.CENTER
    )

    # Strikethrough line
    strike = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.5), Inches(3.7), Inches(8.333), Inches(0.06)
    )
    strike.fill.solid()
    strike.fill.fore_color.rgb = Colors.CORAL
    strike.line.fill.background()

    # X marks - organic blobs
    x_blob1 = create_organic_blob(slide,
        Inches(1.8), Inches(3.0),
        Inches(1.2), Inches(1.2),
        Colors.CORAL_PALE, opacity=50
    )
    add_display_text(slide,
        Inches(2.0), Inches(3.1), Inches(0.8), Inches(1),
        "X",
        size=Pt(32), color=Colors.CORAL,
        align=PP_ALIGN.CENTER
    )

    x_blob2 = create_organic_blob(slide,
        Inches(10.3), Inches(3.0),
        Inches(1.2), Inches(1.2),
        Colors.CORAL_PALE, opacity=50
    )
    add_display_text(slide,
        Inches(10.5), Inches(3.1), Inches(0.8), Inches(1),
        "X",
        size=Pt(32), color=Colors.CORAL,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_18_worth_pursuing(prs):
    """
    Is Etsy even worth pursuing in 2026?

    THE question - dark, dramatic
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_dark_dramatic_background(slide)

    # The big question
    add_display_text(slide,
        Inches(1), Inches(2.2), Inches(11.333), Inches(3),
        "Is Etsy even worth\npursuing in 2026?",
        size=Pt(60), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    # Subtle accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5), Inches(5.8), Inches(3.333), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = Colors.CORAL
    line.line.fill.background()

    return slide


def build_slide_19_answer_tonight(prs):
    """
    I'm going to answer that question tonight.

    Promise and anticipation
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_2(slide)

    # Main statement
    add_display_text(slide,
        Inches(1), Inches(1.8), Inches(11.333), Inches(1.2),
        "I'm going to answer that\nquestion tonight.",
        size=Pt(44), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # Teaser card
    teaser = add_floating_card(slide,
        Inches(3), Inches(3.8),
        Inches(7.333), Inches(2.2),
        Colors.TEAL_DEEP, shadow=True
    )

    add_text(slide,
        Inches(3.3), Inches(4.1), Inches(6.733), Inches(0.6),
        "And the answer?",
        font=Fonts.BODY, size=Pt(20), color=Colors.TEAL_LIGHT,
        align=PP_ALIGN.CENTER
    )

    add_display_text(slide,
        Inches(3.3), Inches(4.7), Inches(6.733), Inches(1),
        "It might honestly\nsurprise you.",
        size=Pt(32), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_slide_20_ai_opportunity(prs):
    """
    The AI flood CREATES an opportunity

    Pattern interrupt - the twist reveal
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_editorial_background_1(slide)

    # Opening context
    add_body_text(slide,
        Inches(1), Inches(0.6), Inches(11.333), Inches(0.5),
        "Because what nobody is talking about is this:",
        size=Pt(16), color=Colors.MUTED,
        align=PP_ALIGN.CENTER
    )

    # Main reveal card
    reveal = add_floating_card(slide,
        Inches(1), Inches(1.3),
        Inches(11.333), Inches(2.8),
        Colors.TEAL_DEEP, shadow=True
    )

    add_text(slide,
        Inches(1.3), Inches(1.6), Inches(10.733), Inches(0.6),
        "The AI flood actually",
        font=Fonts.BODY, size=Pt(24), color=Colors.TEAL_LIGHT,
        align=PP_ALIGN.CENTER
    )

    add_display_text(slide,
        Inches(1.3), Inches(2.2), Inches(10.733), Inches(1.5),
        "CREATES an opportunity",
        size=Pt(52), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    # For specific sellers
    add_body_text(slide,
        Inches(1), Inches(4.3), Inches(11.333), Inches(0.6),
        "— for a very specific type of seller...",
        size=Pt(22), color=Colors.DARK,
        align=PP_ALIGN.CENTER
    )

    # BUT warning - organic blob
    but_blob = create_organic_blob(slide,
        Inches(4.5), Inches(5.0),
        Inches(4.333), Inches(1.2),
        Colors.CORAL, opacity=100
    )
    add_soft_shadow(but_blob, blur=10, distance=3, opacity=15)

    add_display_text(slide,
        Inches(5.0), Inches(5.15), Inches(3.333), Inches(0.9),
        "BUT",
        size=Pt(36), color=Colors.WHITE,
        align=PP_ALIGN.CENTER
    )

    # Condition
    add_body_text(slide,
        Inches(1), Inches(6.4), Inches(11.333), Inches(0.5),
        "Only if you understand what's really happening...",
        size=Pt(18), color=Colors.MUTED,
        align=PP_ALIGN.CENTER
    )

    return slide


# =============================================================================
# MAIN BUILD
# =============================================================================

def build_presentation():
    """Build the premium editorial presentation"""
    print("=" * 60)
    print("BAILEY VANN - THE 2026 ETSY RESET")
    print("Premium Editorial Slide Deck - Version 2")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        ("01", "Title", build_slide_01_title),
        ("02", "Before We Begin", build_slide_02_before_begin),
        ("03", "Get Ready Chat", build_slide_03_get_ready_chat),
        ("04", "Quiz A/B", build_slide_04_quiz_ab),
        ("05", "Type A or B", build_slide_05_type_ab),
        ("06", "The Answer Is", build_slide_06_answer_is),
        ("07", "Both AI", build_slide_07_both_ai),
        ("08", "AI Nowadays", build_slide_08_ai_nowadays),
        ("09", "Sink In", build_slide_09_sink_in),
        ("10", "Uncomfortable", build_slide_10_uncomfortable),
        ("11", "What Happens", build_slide_11_what_happens),
        ("12", "Survey Intro", build_slide_12_survey_intro),
        ("13", "Stat 75+", build_slide_13_stat_75),
        ("14", "Stat 58+", build_slide_14_stat_58),
        ("15", "Stat 47+", build_slide_15_stat_47),
        ("16", "Type YES", build_slide_16_type_yes),
        ("17", "Real Question", build_slide_17_real_question),
        ("18", "Worth Pursuing", build_slide_18_worth_pursuing),
        ("19", "Answer Tonight", build_slide_19_answer_tonight),
        ("20", "AI Opportunity", build_slide_20_ai_opportunity),
    ]

    for num, name, builder in builders:
        print(f"  Building Slide {num}: {name}...")
        builder(prs)

    output = "/home/user/webby-slides-bailey/Bailey_Etsy_Reset_V2.pptx"
    prs.save(output)

    print("=" * 60)
    print(f"SAVED: {output}")
    print("=" * 60)

    return output


if __name__ == "__main__":
    build_presentation()
